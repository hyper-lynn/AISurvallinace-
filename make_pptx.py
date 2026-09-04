import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set slide width and height to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Constants
    BG_COLOR = RGBColor(15, 23, 42)      # Deep Navy Blue #0f172a
    CARD_BG = RGBColor(30, 41, 59)       # Dark Slate #1e293b
    ACCENT_BLUE = RGBColor(59, 130, 246)  # Electric Blue #3b82f6
    ACCENT_CYAN = RGBColor(6, 182, 212)   # Cyan #06b6d4
    TEXT_WHITE = RGBColor(255, 255, 255)  # Pure White
    TEXT_MUTED = RGBColor(148, 163, 184) # Muted Gray #94a3b8

    def set_background(slide):
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        background.fill.solid()
        background.fill.fore_color.rgb = BG_COLOR
        background.line.fill.background() # No border
        return background

    def add_header(slide, title_text, category_text="S-EYE SECURITY PLATFORM"):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
        tf = header_box.text_frame
        tf.word_wrap = True
        
        # Category / Subtitle
        p1 = tf.paragraphs[0]
        p1.text = category_text.upper()
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_CYAN
        
        # Main Title
        p2 = tf.add_paragraph()
        p2.text = title_text
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Front Cover)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)

    # Left content box
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6.5), Inches(5.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "S-EYE INTELLIGENT SURVEILLANCE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p = tf1.add_paragraph()
    p.text = "Real-Time AI Industrial Security"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)

    p = tf1.add_paragraph()
    p.text = "with Media Alert Integration"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(20)

    bullets = [
        "24/7 Autonomous AI Video Stream Analysis",
        "Human Motion Detection & Face Recognition",
        "Instant Telegram Photo & Video Clip Alerts",
        "Zero Operator Fatigue & Proactive Deterrence"
    ]
    for b in bullets:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(6)

    # Embed PNG image on right side if exists
    img_path = r"e:\project\S Eye\pamphlet_design.png"
    if os.path.exists(img_path):
        slide1.shapes.add_picture(img_path, Inches(7.6), Inches(0.8), width=Inches(4.9))

    # -------------------------------------------------------------
    # SLIDE 2: Core AI Vision Technologies
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)
    add_header(slide2, "Core AI Detection & Analytics Engines", "TECHNOLOGY OVERVIEW")

    # Card 1: Human Motion Detection
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = ACCENT_BLUE

    tf_c1 = card1.text_frame
    tf_c1.word_wrap = True
    p = tf_c1.paragraphs[0]
    p.text = "Human Motion Detection Engine"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(14)

    items1 = [
        "Real-Time Stream Processing: Ingests RTSP feeds directly from IP cameras.",
        "Restricted Zone Intrusion: Detects unauthorized entry into sensitive areas.",
        "Off-Hours Surveillance: Automatically monitors facility perimeters at night.",
        "Resource Efficiency: Discards unoccupied frames to optimize server load."
    ]
    for item in items1:
        p = tf_c1.add_paragraph()
        p.text = "✔ " + item
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)

    # Card 2: Face Recognition Engine
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = ACCENT_CYAN

    tf_c2 = card2.text_frame
    tf_c2.word_wrap = True
    p = tf_c2.paragraphs[0]
    p.text = "Face Recognition & Identity Verification"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(14)

    items2 = [
        "Biometric Verification: Compares detected faces against authorized encodings.",
        "Employee vs Stranger Classification: Instantly flags unauthorized strangers.",
        "Encrypted Local Database: Securely stores employee face profiles in SQLite.",
        "Low False-Alarm Rate: High precision detection eliminates false security alerts."
    ]
    for item in items2:
        p = tf_c2.add_paragraph()
        p.text = "✔ " + item
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 3: Media Alert Integration
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)
    add_header(slide3, "Real-Time Multi-Channel Media Alert Integration", "ALERT DISPATCH SYSTEM")

    features = [
        ("Rich Photo Snapshots", "Captures crystal-clear high-definition photo snapshots at the exact moment of security breach.", ACCENT_BLUE),
        ("Video Clip Recording", "Records short video clips of the event to provide immediate visual context to security teams.", ACCENT_CYAN),
        ("Telegram Bot Dispatcher", "Delivers alerts directly to private chats, tactical group channels, and broadcast security HQs.", ACCENT_BLUE),
        ("Mobile & Desktop App Push", "Cross-platform notifications containing timestamp, camera location ID, and threat level status.", ACCENT_CYAN)
    ]

    for i, (title, desc, color) in enumerate(features):
        x = Inches(0.8 + (i % 2) * 6.1)
        y = Inches(1.8 + (i // 2) * 2.5)
        
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------
    # SLIDE 4: System Architecture Flowchart
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)
    add_header(slide4, "End-to-End System Workflow Architecture", "SYSTEM ARCHITECTURE")

    flow_path = r"e:\project\S Eye\project_flowchart.png"
    if os.path.exists(flow_path):
        slide4.shapes.add_picture(flow_path, Inches(0.8), Inches(1.6), width=Inches(11.733))
    else:
        tb = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.733), Inches(4.0))
        tb.text_frame.text = "Flowchart image project_flowchart.png not found."

    # -------------------------------------------------------------
    # SLIDE 5: Key Business & Security Benefits
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)
    add_header(slide5, "Key Security, Financial & Operational Benefits", "BUSINESS IMPACT")

    benefits = [
        ("Proactive Security Deterrence", "Shifts security operations from post-incident review to real-time incident prevention."),
        ("Zero Operator Fatigue", "AI monitors all cameras 24/7 without fatigue, eliminating human monitoring errors."),
        ("Cost Efficiency & Scaling", "Integrates seamlessly with existing CCTV hardware, avoiding expensive infrastructure upgrades."),
        ("Verifiable Audit Logging", "Maintains complete digital logs with images, timestamps, and face IDs for compliance audits.")
    ]

    for i, (b_title, b_desc) in enumerate(benefits):
        y = Inches(1.7 + i * 1.3)
        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.733), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_BLUE if i % 2 == 0 else ACCENT_CYAN

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = b_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN if i % 2 == 0 else ACCENT_BLUE
        
        p = tf.add_paragraph()
        p.text = b_desc
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------
    # SLIDE 6: Technical Specs & Contact Information
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6)
    add_header(slide6, "Technical Specifications & Contact Details", "SPECIFICATIONS & CONTACT")

    # Specs Card
    spec_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    spec_card.fill.solid()
    spec_card.fill.fore_color.rgb = CARD_BG
    spec_card.line.color.rgb = ACCENT_BLUE
    
    tf_s = spec_card.text_frame
    tf_s.word_wrap = True
    p = tf_s.paragraphs[0]
    p.text = "Technical Specifications"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(10)

    specs = [
        "Camera Input: RTSP Streams, IP Cameras, NVR/DVR",
        "AI Vision Engine: Human Motion & Face Recognition",
        "Alert Destination: Telegram Bot API, Mobile/Desktop App",
        "Database Engine: Encrypted SQLite Persistence Layer",
        "Deployment: Local Server / Cloud Hybrid Architecture"
    ]
    for s in specs:
        p = tf_s.add_paragraph()
        p.text = "• " + s
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(8)

    # Contact Card
    contact_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    contact_card.fill.solid()
    contact_card.fill.fore_color.rgb = CARD_BG
    contact_card.line.color.rgb = ACCENT_CYAN

    tf_ct = contact_card.text_frame
    tf_ct.word_wrap = True
    p = tf_ct.paragraphs[0]
    p.text = "Get Started with S-Eye"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(10)

    contacts = [
        "Product Name: S-Eye Security Platform",
        "System Status: Production Ready / Custom Integration",
        "Email Support: support@s-eye-security.com",
        "Official Website: www.s-eye-security.com",
        "Request Demo: Contact Security Integration Team"
    ]
    for c in contacts:
        p = tf_ct.add_paragraph()
        p.text = "• " + c
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(8)

    # Save presentation
    output_path = r"e:\project\S Eye\S_Eye_Security_Pamphlet_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    create_presentation()
